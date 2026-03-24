#!/usr/bin/env python3
"""
Generate a clean, professional medical PowerPoint on Nail Clubbing.
Designed for Physiotherapy PG students. Copyright-friendly content.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Clean Medical Palette ──
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xFA, 0xFA, 0xFC)
NAVY        = RGBColor(0x1A, 0x2B, 0x4B)
SLATE       = RGBColor(0x33, 0x44, 0x55)
BODY_TEXT   = RGBColor(0x3D, 0x4F, 0x5F)
MUTED       = RGBColor(0x6B, 0x7B, 0x8D)
TEAL        = RGBColor(0x00, 0x7B, 0x83)
LIGHT_TEAL  = RGBColor(0xE0, 0xF2, 0xF1)
BORDER      = RGBColor(0xD0, 0xD8, 0xE0)
GOLD        = RGBColor(0xC8, 0x96, 0x2E)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

TOTAL_SLIDES = 11

# ── Helpers ──

def solid_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def rrect(slide, l, t, w, h, fill, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def txt(slide, l, t, w, h, text, sz=16, color=BODY_TEXT, bold=False,
        align=PP_ALIGN.LEFT, font="Calibri", spacing=1.15):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.space_after = Pt(0)
    if spacing != 1.0:
        p.line_spacing = Pt(sz * spacing)
    return tb

def bullets(slide, l, t, w, h, items, sz=14, color=BODY_TEXT, sp=6):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(sp)
    return tb

def multi_text(slide, l, t, w, h, lines, sz=14, color=BODY_TEXT, sp=6):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(sp)
    return tb

def header_bar(slide, title):
    rect(slide, Inches(0), Inches(0), W, Inches(1.05), NAVY)
    rect(slide, Inches(0), Inches(1.05), W, Inches(0.04), TEAL)
    txt(slide, Inches(0.8), Inches(0.22), Inches(11), Inches(0.6),
        title, sz=28, color=WHITE, bold=True)

def slide_num(slide, n):
    txt(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.35),
        f"{n} / {TOTAL_SLIDES}", sz=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════
# 1 — TITLE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, NAVY)

rect(s, Inches(0), Inches(0), W, Inches(0.035), TEAL)

txt(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.0),
    "Nail Clubbing", sz=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
    font="Calibri")

rect(s, Inches(5.8), Inches(3.2), Inches(1.7), Inches(0.03), GOLD)

txt(s, Inches(1), Inches(3.5), Inches(11.3), Inches(0.6),
    "Definition  •  Mechanism  •  Grades  •  Tests  •  Causes",
    sz=18, color=RGBColor(0xA0, 0xB0, 0xC0), align=PP_ALIGN.CENTER)

txt(s, Inches(1), Inches(5.2), Inches(11.3), Inches(0.4),
    "Presented by: [Your Name]", sz=16, color=RGBColor(0x8A, 0x9A, 0xAA),
    align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.65), Inches(11.3), Inches(0.4),
    "PG Student, Physiotherapy  •  [College Name]  •  [Date]",
    sz=13, color=RGBColor(0x70, 0x80, 0x90), align=PP_ALIGN.CENTER)

rect(s, Inches(0), Inches(7.465), W, Inches(0.035), TEAL)
slide_num(s, 1)


# ═══════════════════════════════════════════════════════════════
# 2 — DEFINITION
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, OFF_WHITE)
header_bar(s, "Definition")

# Definition card
rrect(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.8), WHITE, border=TEAL)
rect(s, Inches(0.8), Inches(1.5), Inches(0.06), Inches(1.8), TEAL)

txt(s, Inches(1.15), Inches(1.65), Inches(11.0), Inches(1.5),
    "Digital clubbing (Hippocratic fingers) is a clinical sign characterised by "
    "bulbous, uniform swelling of the soft tissue of the terminal phalanx of a digit, "
    "with loss of the normal angle (Lovibond angle) between the nail plate and the "
    "proximal nail fold.",
    sz=17, color=SLATE, spacing=1.5)

# Synonym & first description
txt(s, Inches(0.8), Inches(3.7), Inches(5), Inches(0.35),
    "Also known as:", sz=13, color=MUTED, bold=True)
txt(s, Inches(0.8), Inches(4.05), Inches(6), Inches(0.6),
    "Hippocratic fingers, Watch-glass nails, Drumstick fingers",
    sz=14, color=BODY_TEXT)

txt(s, Inches(7.0), Inches(3.7), Inches(5), Inches(0.35),
    "First described:", sz=13, color=MUTED, bold=True)
txt(s, Inches(7.0), Inches(4.05), Inches(6), Inches(0.6),
    "Hippocrates (~400 BCE) in patients with empyema",
    sz=14, color=BODY_TEXT)

# Key features
txt(s, Inches(0.8), Inches(5.0), Inches(4), Inches(0.35),
    "Key Clinical Features", sz=15, color=TEAL, bold=True)

bullets(s, Inches(0.8), Inches(5.4), Inches(11.5), Inches(2.0), [
    "Increased nail curvature in both longitudinal and transverse axes",
    "Sponginess (fluctuation) of the nail bed on palpation",
    "Loss of the normal ~160° Lovibond angle (angle becomes ≥ 180°)",
    "Distal phalangeal depth (DPD) / interphalangeal depth (IPD) ratio > 1.0",
    "May be unilateral or bilateral; hereditary or acquired",
], sz=14, color=BODY_TEXT, sp=5)

slide_num(s, 2)


# ═══════════════════════════════════════════════════════════════
# 3 — PATHOPHYSIOLOGY / MECHANISM
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, OFF_WHITE)
header_bar(s, "Pathophysiology & Mechanism")

txt(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.6),
    "The exact mechanism remains debated. Three principal theories are recognised:",
    sz=15, color=MUTED, spacing=1.3)

theories = [
    ("VEGF / Hypoxia Theory", TEAL, [
        "Chronic hypoxia → increased VEGF release from distal tissues",
        "VEGF promotes angiogenesis and vascular permeability",
        "Connective tissue hyperplasia in the nail bed",
        "Results in soft-tissue swelling of the fingertip",
    ]),
    ("Megakaryocyte / Platelet Theory", RGBColor(0x5C, 0x35, 0x85), [
        "Megakaryocyte fragments normally fragmented in pulmonary capillaries",
        "In cardiopulmonary disease → fragments bypass the lungs",
        "Lodge in distal digital capillaries",
        "Release PDGF & VEGF → local tissue proliferation",
    ]),
    ("Neural & Genetic Factors", RGBColor(0xBF, 0x36, 0x0C), [
        "Vagal-mediated mechanism (clubbing may resolve after vagotomy)",
        "PGE₂ overproduction in hypertrophic osteoarthropathy",
        "HPGD gene mutation in hereditary forms (pachydermoperiostosis)",
        "Unilateral clubbing suggests local vascular aetiology",
    ]),
]

for i, (title, color, pts) in enumerate(theories):
    x = Inches(0.5) + Inches(i * 4.15)
    rrect(s, x, Inches(2.0), Inches(3.95), Inches(5.1), WHITE, border=BORDER)
    rect(s, x, Inches(2.0), Inches(3.95), Inches(0.5), color)
    txt(s, x + Inches(0.2), Inches(2.07), Inches(3.5), Inches(0.4),
        title, sz=13, color=WHITE, bold=True)
    bullets(s, x + Inches(0.2), Inches(2.7), Inches(3.5), Inches(4.0),
            pts, sz=13, color=BODY_TEXT, sp=8)

slide_num(s, 3)


# ═══════════════════════════════════════════════════════════════
# 4 — CAUSES & TYPES (COMBINED)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, OFF_WHITE)
header_bar(s, "Causes & Classification")

# Types row at top
txt(s, Inches(0.8), Inches(1.3), Inches(5), Inches(0.35),
    "Classification of Clubbing", sz=15, color=TEAL, bold=True)

type_data = [
    ("Acquired", "Secondary to underlying disease\n(most common presentation)"),
    ("Hereditary / Idiopathic", "Primary — no underlying disease\n(pachydermoperiostosis, familial)"),
    ("Unilateral", "Suggests local vascular cause\n(AV fistula, Pancoast tumour)"),
]

for i, (label, desc) in enumerate(type_data):
    x = Inches(0.5) + Inches(i * 4.15)
    rrect(s, x, Inches(1.7), Inches(3.95), Inches(1.1), WHITE, border=BORDER)
    txt(s, x + Inches(0.2), Inches(1.77), Inches(3.5), Inches(0.3),
        label, sz=14, color=NAVY, bold=True)
    txt(s, x + Inches(0.2), Inches(2.1), Inches(3.5), Inches(0.65),
        desc, sz=12, color=MUTED)

# Causes — 3 columns below
cause_groups = [
    ("Pulmonary", TEAL, [
        "Bronchogenic carcinoma",
        "Bronchiectasis",
        "Cystic fibrosis",
        "Idiopathic pulmonary fibrosis",
        "Lung abscess / empyema",
        "Mesothelioma",
        "Pulmonary AV malformations",
    ]),
    ("Cardiac", RGBColor(0xC6, 0x28, 0x28), [
        "Cyanotic congenital heart disease",
        "Infective endocarditis (subacute)",
        "Atrial myxoma",
        "Arteriovenous fistulae",
    ]),
    ("Other Systemic", RGBColor(0x55, 0x6B, 0x2F), [
        "Inflammatory bowel disease",
        "Hepatic cirrhosis",
        "Coeliac disease",
        "Graves' disease (thyroid acropachy)",
        "Chronic infections (TB, HIV)",
    ]),
]

for i, (title, color, pts) in enumerate(cause_groups):
    x = Inches(0.5) + Inches(i * 4.15)
    rrect(s, x, Inches(3.15), Inches(3.95), Inches(4.1), WHITE, border=BORDER)
    rect(s, x, Inches(3.15), Inches(3.95), Inches(0.45), color)
    txt(s, x + Inches(0.2), Inches(3.2), Inches(3.5), Inches(0.35),
        title, sz=13, color=WHITE, bold=True)
    bullets(s, x + Inches(0.15), Inches(3.75), Inches(3.6), Inches(3.4),
            pts, sz=13, color=BODY_TEXT, sp=6)

slide_num(s, 4)


# ═══════════════════════════════════════════════════════════════
# 5 — GRADES OF CLUBBING
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, OFF_WHITE)
header_bar(s, "Grades of Clubbing")

grade_colors = [
    RGBColor(0x2E, 0x7D, 0x32),
    RGBColor(0xEF, 0x6C, 0x00),
    RGBColor(0xD3, 0x2F, 0x2F),
    RGBColor(0x6A, 0x1B, 0x9A),
]

grades = [
    ("Grade I", "Fluctuation & Softening", [
        "Nail bed feels spongy on palpation",
        "Loss of normal firm attachment",
        "Lovibond angle still ~160°",
        "Earliest detectable sign",
    ]),
    ("Grade II", "Loss of Lovibond Angle", [
        "Lovibond angle obliterated (≥ 180°)",
        "Profile ratio DPD/IPD > 1.0",
        "Nail appears to 'float'",
        "Schamroth window test becomes positive",
    ]),
    ("Grade III", "Increased Curvature", [
        "Accentuated convexity in all planes",
        "Parrot-beak / watch-glass appearance",
        "Bulbous enlargement of fingertip",
        "Both longitudinal & transverse curvature ↑",
    ]),
    ("Grade IV", "Drumstick / HOA", [
        "Entire distal phalanx swollen",
        "Classic 'drumstick' shape",
        "May have periostitis (wrist, ankle)",
        "Hypertrophic osteoarthropathy (HOA)",
    ]),
]

for i, (grade, subtitle, pts) in enumerate(grades):
    x = Inches(0.35) + Inches(i * 3.2)
    gc = grade_colors[i]
    rrect(s, x, Inches(1.4), Inches(3.0), Inches(5.7), WHITE, border=BORDER)
    # Header
    rect(s, x, Inches(1.4), Inches(3.0), Inches(0.8), gc)
    txt(s, x + Inches(0.15), Inches(1.43), Inches(2.7), Inches(0.4),
        grade, sz=18, color=WHITE, bold=True)
    txt(s, x + Inches(0.15), Inches(1.82), Inches(2.7), Inches(0.3),
        subtitle, sz=11, color=RGBColor(0xFF, 0xFF, 0xDD))
    # Points
    bullets(s, x + Inches(0.12), Inches(2.45), Inches(2.7), Inches(4.4),
            pts, sz=13, color=BODY_TEXT, sp=8)

slide_num(s, 5)


# ═══════════════════════════════════════════════════════════════
# 6 — CLINICAL TESTS OVERVIEW
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, OFF_WHITE)
header_bar(s, "Clinical Assessment & Tests")

tests = [
    ("Schamroth Window Test",
     "Oppose dorsal surfaces of terminal phalanges (nail-to-nail). "
     "A diamond-shaped window is normally visible; obliterated in clubbing.",
     "Quick bedside test  •  Sensitivity ~93%  •  Specificity ~95%"),
    ("Lovibond Angle (Profile Sign)",
     "Measure the angle between the nail plate and proximal nail fold. "
     "Normal < 180° (~160°). Clubbing ≥ 180°.",
     "Objective measurement  •  Can be quantified with digital photography"),
    ("Phalangeal Depth Ratio",
     "DPD (distal phalangeal depth) ÷ IPD (interphalangeal depth). "
     "Normal < 1.0. Clubbing ≥ 1.0.",
     "Most objective and reproducible measure"),
    ("Nail Bed Fluctuation",
     "Press on the nail bed — in early clubbing the nail feels 'floating' or "
     "spongy due to soft-tissue oedema and increased vascularity.",
     "Earliest sign  •  Subjective but clinically useful"),
]

for i, (title, desc, note) in enumerate(tests):
    y = Inches(1.35) + Inches(i * 1.5)
    rrect(s, Inches(0.5), y, Inches(12.3), Inches(1.35), WHITE, border=BORDER)
    rect(s, Inches(0.5), y, Inches(0.05), Inches(1.35), TEAL)
    txt(s, Inches(0.85), y + Pt(6), Inches(11.5), Inches(0.3),
        title, sz=15, color=NAVY, bold=True)
    txt(s, Inches(0.85), y + Pt(28), Inches(11.5), Inches(0.55),
        desc, sz=13, color=BODY_TEXT, spacing=1.3)
    txt(s, Inches(0.85), y + Pt(68), Inches(11.5), Inches(0.25),
        note, sz=11, color=TEAL)

slide_num(s, 6)


# ═══════════════════════════════════════════════════════════════
# 7 — SCHAMROTH TEST DETAIL
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, OFF_WHITE)
header_bar(s, "Schamroth Window Test — Detail")

# Normal column
rrect(s, Inches(0.5), Inches(1.4), Inches(5.9), Inches(5.7), WHITE, border=RGBColor(0x2E, 0x7D, 0x32))
rect(s, Inches(0.5), Inches(1.4), Inches(5.9), Inches(0.5), RGBColor(0x2E, 0x7D, 0x32))
txt(s, Inches(0.7), Inches(1.45), Inches(5.5), Inches(0.4),
    "Normal (Negative)", sz=15, color=WHITE, bold=True)

multi_text(s, Inches(0.8), Inches(2.15), Inches(5.3), Inches(4.5), [
    "• Oppose terminal phalanges of the same finger nail-to-nail",
    "",
    "• A small diamond-shaped 'window' of light is visible",
    "  between the nail bases",
    "",
    "• This gap exists because Lovibond angle < 180°",
    "",
    "• Described by Leo Schamroth (1976) after observing loss",
    "  of this window during his own endocarditis",
], sz=14, color=BODY_TEXT, sp=3)

# Positive column
rrect(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(5.7), WHITE, border=RGBColor(0xC6, 0x28, 0x28))
rect(s, Inches(6.9), Inches(1.4), Inches(5.9), Inches(0.5), RGBColor(0xC6, 0x28, 0x28))
txt(s, Inches(7.1), Inches(1.45), Inches(5.5), Inches(0.4),
    "Clubbing (Positive)", sz=15, color=WHITE, bold=True)

multi_text(s, Inches(7.2), Inches(2.15), Inches(5.3), Inches(4.5), [
    "• The diamond-shaped window is ABSENT (obliterated)",
    "",
    "• Lovibond angle ≥ 180° → nail folds flush against each",
    "  other → no gap remains",
    "",
    "• Interpretation:",
    "    Window present → Negative (no clubbing)",
    "    Window absent → Positive (clubbing present)",
    "",
    "• Limitation: less reliable in early Grade I clubbing;",
    "  combine with phalangeal depth ratio for accuracy",
], sz=14, color=BODY_TEXT, sp=3)

slide_num(s, 7)


# ═══════════════════════════════════════════════════════════════
# 8 — DIFFERENTIAL DIAGNOSIS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, OFF_WHITE)
header_bar(s, "Differential Diagnosis")

txt(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.4),
    "Conditions that may mimic clubbing:", sz=14, color=MUTED)

diffs = [
    ("Pseudoclubbing",
     "Nail curvature ↑ but NO soft-tissue increase. DPD/IPD < 1.0. Seen in renal osteodystrophy."),
    ("Koilonychia",
     "Spoon-shaped (concave) nails — opposite of clubbing. Associated with iron-deficiency anaemia."),
    ("Paronychia",
     "Periungual tissue infection → localised swelling. Usually painful and unilateral."),
    ("Felons / Whitlows",
     "Fingertip pulp swelling from infection. Tender, warm, erythematous — unlike painless clubbing."),
    ("Onycholysis",
     "Nail plate separates from nail bed. Can mimic the 'floating nail' but different cause."),
]

for i, (cond, desc) in enumerate(diffs):
    y = Inches(1.85) + Inches(i * 1.05)
    rrect(s, Inches(0.5), y, Inches(12.3), Inches(0.9), WHITE, border=BORDER)
    rect(s, Inches(0.5), y, Inches(0.05), Inches(0.9), TEAL)
    txt(s, Inches(0.85), y + Pt(5), Inches(3), Inches(0.3),
        cond, sz=14, color=NAVY, bold=True)
    txt(s, Inches(0.85), y + Pt(26), Inches(11.5), Inches(0.5),
        desc, sz=13, color=BODY_TEXT, spacing=1.2)

slide_num(s, 8)


# ═══════════════════════════════════════════════════════════════
# 9 — PHYSIOTHERAPY RELEVANCE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, OFF_WHITE)
header_bar(s, "Relevance to Physiotherapy Practice")

sections = [
    ("Clinical Assessment", TEAL, [
        "Clubbing is a red flag for underlying cardiopulmonary pathology",
        "Part of routine hand inspection in every cardiorespiratory examination",
        "Progressive clubbing → refer for medical investigation",
        "Regression may indicate response to treatment",
    ]),
    ("Cardiopulmonary Rehabilitation", RGBColor(0xC6, 0x28, 0x28), [
        "Patients often have chronic lung disease (COPD, IPF, bronchiectasis)",
        "Tailor exercise prescription to underlying diagnosis",
        "Monitor SpO₂ — hypoxia may be present even at rest",
        "Airway clearance techniques in bronchiectasis / CF patients",
    ]),
    ("Patient Education", RGBColor(0xEF, 0x6C, 0x00), [
        "Explain clubbing as a systemic sign — not a local nail problem",
        "Motivate adherence to pulmonary rehabilitation programmes",
        "Teach self-monitoring (oxygen saturation, symptom diaries)",
        "Encourage smoking cessation where relevant",
    ]),
]

for i, (title, color, pts) in enumerate(sections):
    x = Inches(0.35) + Inches(i * 4.2)
    rrect(s, x, Inches(1.4), Inches(4.0), Inches(5.7), WHITE, border=BORDER)
    rect(s, x, Inches(1.4), Inches(4.0), Inches(0.5), color)
    txt(s, x + Inches(0.2), Inches(1.45), Inches(3.6), Inches(0.4),
        title, sz=14, color=WHITE, bold=True)
    bullets(s, x + Inches(0.15), Inches(2.1), Inches(3.7), Inches(4.8),
            pts, sz=13, color=BODY_TEXT, sp=8)

slide_num(s, 9)


# ═══════════════════════════════════════════════════════════════
# 10 — REFERENCES
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, OFF_WHITE)
header_bar(s, "References")

txt(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.35),
    "Textbooks — Cardiopulmonary & General Medicine", sz=15, color=TEAL, bold=True)

refs1 = [
    "1.  Hillegass, E. (2022). Essentials of Cardiopulmonary Physical Therapy. 5th ed. Elsevier.",
    "2.  Watchie, J. (2010). Cardiovascular and Pulmonary Physical Therapy. 2nd ed. Saunders.",
    "3.  Pryor, J.A. & Prasad, A.S. (2008). Physiotherapy for Respiratory and Cardiac Problems. 4th ed. Churchill Livingstone.",
    "4.  Frownfelter, D. & Dean, E. (2012). Cardiovascular and Pulmonary Physical Therapy. 5th ed. Mosby.",
    "5.  West, J.B. (2021). West's Respiratory Physiology: The Essentials. 11th ed. Wolters Kluwer.",
    "6.  Kumar, P. & Clark, M. (2021). Kumar & Clark's Clinical Medicine. 10th ed. Elsevier.",
]

multi_text(s, Inches(0.8), Inches(1.75), Inches(11.5), Inches(2.2),
           refs1, sz=13, color=BODY_TEXT, sp=7)

txt(s, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.35),
    "Journal Articles & Seminal Papers", sz=15, color=TEAL, bold=True)

refs2 = [
    "7.  Schamroth, L. (1976). Personal experience. S Afr Med J, 50(9), 297–300.",
    "8.  Spicknall, K.E. & Zirwas, M.J. (2009). Clubbing: an update. J Am Acad Dermatol, 60(6), 1073–1082.",
    "9.  Dickinson, C.J. (1993). The aetiology of clubbing and HOA. Eur J Clin Invest, 23(6), 330–338.",
    "10. Lovibond, J.L. (1938). Diagnosis of clubbed fingers. The Lancet, 231(5979), 363–364.",
    "11. Callemeyn, J. et al. (2016). Clubbing and HOA. Clin Rheumatol, 35, 2575–2581.",
]

multi_text(s, Inches(0.8), Inches(4.65), Inches(11.5), Inches(2.2),
           refs2, sz=13, color=MUTED, sp=7)

slide_num(s, 10)


# ═══════════════════════════════════════════════════════════════
# 11 — THANK YOU
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, NAVY)

rect(s, Inches(0), Inches(0), W, Inches(0.035), TEAL)

txt(s, Inches(1), Inches(2.8), Inches(11.3), Inches(0.8),
    "Thank You", sz=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

rect(s, Inches(5.8), Inches(3.8), Inches(1.7), Inches(0.03), GOLD)

txt(s, Inches(1), Inches(4.1), Inches(11.3), Inches(0.5),
    "Questions & Discussion", sz=20, color=RGBColor(0xA0, 0xB0, 0xC0),
    align=PP_ALIGN.CENTER)

txt(s, Inches(1), Inches(5.5), Inches(11.3), Inches(0.4),
    "[Your Name]  •  [Email]",
    sz=14, color=RGBColor(0x70, 0x80, 0x90), align=PP_ALIGN.CENTER)

rect(s, Inches(0), Inches(7.465), W, Inches(0.035), TEAL)
slide_num(s, 11)


# ═══════════════════════════ SAVE ═══════════════════════════
output_path = os.path.expanduser("~/Documents/Nail_Clubbing_Presentation.pptx")
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {TOTAL_SLIDES}")
